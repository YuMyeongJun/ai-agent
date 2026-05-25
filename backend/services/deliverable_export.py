from __future__ import annotations

import re
from pathlib import Path
from xml.sax.saxutils import escape

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas

_FONT_NAME: str | None = None


def _resolve_pdf_font() -> str:
    global _FONT_NAME
    if _FONT_NAME:
        return _FONT_NAME

    candidates = [
        "/System/Library/Fonts/Supplemental/AppleGothic.ttf",
        "/Library/Fonts/Arial Unicode.ttf",
        "/usr/share/fonts/truetype/nanum/NanumGothic.ttf",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    ]
    for path in candidates:
        if Path(path).exists():
            try:
                pdfmetrics.registerFont(TTFont("OfficeFont", path))
                _FONT_NAME = "OfficeFont"
                return _FONT_NAME
            except Exception:
                continue

    _FONT_NAME = "Helvetica"
    return _FONT_NAME


def _split_sections(text: str) -> list[tuple[str, str]]:
    sections: list[tuple[str, str]] = []
    current_title = "Overview"
    current_lines: list[str] = []

    for line in text.splitlines():
        if line.startswith("# "):
            if current_lines:
                sections.append((current_title, "\n".join(current_lines).strip()))
            current_title = line[2:].strip()
            current_lines = []
        elif line.startswith("## "):
            if current_lines:
                sections.append((current_title, "\n".join(current_lines).strip()))
            current_title = line[3:].strip()
            current_lines = []
        else:
            current_lines.append(line)

    if current_lines or not sections:
        sections.append((current_title, "\n".join(current_lines).strip() or text.strip()))

    return sections[:12]


def _wrap_pdf_line(c: canvas.Canvas, text: str, x: float, y: float, max_width: float, font: str, size: int) -> float:
    c.setFont(font, size)
    words = text.split()
    line = ""
    for word in words:
        trial = f"{line} {word}".strip()
        if c.stringWidth(trial, font, size) <= max_width:
            line = trial
        else:
            if line:
                c.drawString(x, y, line[:120])
                y -= size + 4
            line = word
    if line:
        c.drawString(x, y, line[:120])
        y -= size + 4
    return y


def export_markdown_pdf(path: Path, title: str, markdown: str) -> None:
    font = _resolve_pdf_font()
    c = canvas.Canvas(str(path), pagesize=A4)
    width, height = A4
    margin = 20 * mm
    y = height - margin

    c.setFont(font, 18)
    c.drawString(margin, y, title[:60])
    y -= 28

    for section_title, body in _split_sections(markdown):
        if y < margin + 40:
            c.showPage()
            y = height - margin

        c.setFont(font, 13)
        c.drawString(margin, y, section_title[:80])
        y -= 18

        for paragraph in body.split("\n\n"):
            paragraph = paragraph.strip()
            if not paragraph:
                continue
            paragraph = re.sub(r"^[-*]\s+", "• ", paragraph, flags=re.MULTILINE)
            y = _wrap_pdf_line(c, paragraph.replace("\n", " "), margin, y, width - 2 * margin, font, 10)
            y -= 6
            if y < margin + 30:
                c.showPage()
                y = height - margin

    c.save()


def export_presentation_pptx(path: Path, title: str, markdown: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title[:80]
    if slide.placeholders[1]:
        slide.placeholders[1].text = "AI Virtual Office — 기획 발표자료"

    bullet_layout = prs.slide_layouts[1]
    for section_title, body in _split_sections(markdown):
        slide = prs.slides.add_slide(bullet_layout)
        slide.shapes.title.text = section_title[:80]
        tf = slide.placeholders[1].text_frame
        tf.clear()
        lines = [ln.strip() for ln in body.splitlines() if ln.strip()]
        if not lines:
            lines = ["(내용 없음)"]
        for index, line in enumerate(lines[:8]):
            line = re.sub(r"^[-*#]+\s*", "", line)
            if index == 0:
                tf.text = line[:200]
            else:
                tf.add_paragraph().text = line[:200]

    prs.save(str(path))


def export_ui_mockup_html(path: Path, mission: str, design_notes: str) -> None:
    sections = _split_sections(design_notes)
    feature_cards = ""
    for title, body in sections[1:4]:
        lines = [ln.strip("-• ") for ln in body.splitlines() if ln.strip()][:2]
        feature_cards += f"""
        <article class="card">
          <h3>{escape(title[:40])}</h3>
          <p>{escape(" · ".join(lines)[:120] or "Feature highlight")}</p>
        </article>"""

    hero_lines = sections[0][1].splitlines() if sections else []
    hero = hero_lines[0][:120] if hero_lines and hero_lines[0].strip() else mission
    html = f"""<!DOCTYPE html>
<html lang="ko">
<head>
  <meta charset="utf-8"/>
  <meta name="viewport" content="width=device-width, initial-scale=1"/>
  <title>{escape(mission[:60])} — UI Mockup</title>
  <style>
    * {{ box-sizing: border-box; margin: 0; padding: 0; }}
    body {{ font-family: system-ui, sans-serif; background: #0f1117; color: #e8e8f0; }}
    .nav {{ display: flex; justify-content: space-between; padding: 16px 32px; border-bottom: 1px solid #2a2a3a; }}
    .logo {{ font-weight: 700; color: #7ba87b; }}
    .hero {{ display: grid; grid-template-columns: 1fr 1fr; gap: 32px; padding: 48px 32px; max-width: 1100px; margin: 0 auto; }}
    .hero h1 {{ font-size: 2rem; margin-bottom: 12px; }}
    .hero p {{ color: #9898b0; line-height: 1.6; margin-bottom: 20px; }}
    .cta {{ display: inline-block; padding: 12px 24px; background: #7ba87b; color: #fff; border-radius: 8px; text-decoration: none; }}
    .preview {{ background: linear-gradient(135deg,#1a1a2e,#2d2d44); border-radius: 16px; min-height: 220px; border: 1px solid #3a3a50; }}
    .grid {{ display: grid; grid-template-columns: repeat(3,1fr); gap: 16px; max-width: 1100px; margin: 0 auto; padding: 0 32px 48px; }}
    .card {{ background: #1a1a2e; border: 1px solid #2a2a3a; border-radius: 12px; padding: 20px; }}
    .card h3 {{ margin-bottom: 8px; font-size: 1rem; }}
    .card p {{ color: #9898b0; font-size: 0.9rem; }}
    @media (max-width: 768px) {{ .hero, .grid {{ grid-template-columns: 1fr; }} }}
  </style>
</head>
<body>
  <header class="nav">
    <div class="logo">Brand</div>
    <nav>Features · Pricing · Contact</nav>
  </header>
  <section class="hero">
    <div>
      <h1>{escape(mission[:80])}</h1>
      <p>{escape(hero)}</p>
      <a class="cta" href="#">Get Started</a>
    </div>
    <div class="preview"></div>
  </section>
  <section class="grid">{feature_cards or '<article class="card"><h3>Feature</h3><p>Design placeholder</p></article>'}</section>
</body>
</html>"""
    path.write_text(html, encoding="utf-8")


def export_wireframe_svg(path: Path, mission: str, notes: str) -> None:
    blocks = [
        (40, 40, 720, 48, "Header / Navigation"),
        (40, 108, 340, 180, "Hero copy + CTA"),
        (400, 108, 360, 180, "Hero visual"),
        (40, 308, 220, 120, "Feature A"),
        (280, 308, 220, 120, "Feature B"),
        (520, 308, 240, 120, "Feature C"),
        (40, 448, 720, 80, "Footer / Contact"),
    ]
    parts = [
        f'<svg xmlns="http://www.w3.org/2000/svg" width="800" height="540" viewBox="0 0 800 540">',
        f'<rect width="800" height="540" fill="#f5f5f8"/>',
        f'<text x="40" y="24" font-family="system-ui,sans-serif" font-size="14" fill="#444">{escape(mission[:50])} — Wireframe</text>',
    ]
    for x, y, w, h, label in blocks:
        parts.append(f'<rect x="{x}" y="{y}" width="{w}" height="{h}" fill="#fff" stroke="#888" stroke-width="2" rx="6"/>')
        parts.append(
            f'<text x="{x + 12}" y="{y + 24}" font-family="system-ui,sans-serif" font-size="12" fill="#666">{escape(label)}</text>'
        )
    parts.append("</svg>")
    path.write_text("\n".join(parts), encoding="utf-8")
